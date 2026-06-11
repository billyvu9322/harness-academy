# -*- coding: utf-8 -*-
"""Rebuild seminar deck into Anthropic editorial template.
Operates in-place on the already-recolored pptx so the 3 complex diagram
slides (18,19,20) are preserved untouched; all other slides are rebuilt.
Text strings are pulled from deck.json (exact) so nothing is retyped."""
import json
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

SRC = "_base_recolored.pptx"
OUT = "Senimar-Harness-Engineering.pptx"

# ---- palette ----
BG        = "EDEBE4"   # warm cream background
INK       = "1A1A1A"   # near-black / dark card
CARD_L    = "E3E0D7"   # light card (slightly darker cream)
ONDARK    = "F4F0E8"   # light text on dark card
MUTED     = "6E6A60"   # muted label / subtitle on light
BODY      = "33322E"   # body text on light
MUTEDDARK = "B7B2A6"   # muted text on dark card
DARKBG    = "26231E"   # warm charcoal slide background (title/divider/demo slides)
RULELT    = "C9C4B8"   # light hairline rule on dark bg
SERIF = "Georgia"
SANS  = "Inter"
RGB = RGBColor.from_string

deck = json.load(open("deck.json", encoding="utf-8"))
def slide_paras(n):
    """ordered list of paragraph-text lists for each text box on slide n (1-based)"""
    out = []
    for sh in deck[n-1]["shapes"]:
        if "p" in sh:
            out.append([p["x"] for p in sh["p"]])
    return out

prs = Presentation(SRC)
PINK = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

# update master background to warm cream
for m in prs.slide_masters:
    for el in m.element.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"):
        if (el.get("val") or "").upper() in ("F5F0E8",):
            el.set("val", BG)

# strip every slide's own <p:bg> (originals carried dark/light per-slide fills)
# so the cream master background governs uniformly across all 22 slides
for sl in prs.slides:
    for bgel in sl.element.findall(PINK + "cSld/" + PINK + "bg"):
        bgel.getparent().remove(bgel)

# ---------- helpers ----------
def clear_keep_pics(slide):
    # strip orphan-prone animation/transition nodes that reference shape spids
    sld = slide.element
    for tag in ("timing", "transition"):
        for el in sld.findall(PINK + tag):
            sld.remove(el)
    for sh in list(slide.shapes):
        if sh.shape_type != 13:  # keep PICTURE
            sh._element.getparent().remove(sh._element)

def get_pic(slide):
    for sh in slide.shapes:
        if sh.shape_type == 13:
            return sh
    return None

def box(slide, l, t, w, h, anchor=None, wrap=True):
    tb = slide.shapes.add_textbox(In(l), In(t), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    return tb, tf

def para(tf, text, size, color, font=SANS, bold=False, first=False,
         before=0, after=0, line=None, align=None, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if before: p.space_before = Pt(before)
    if after:  p.space_after = Pt(after)
    if line:   p.line_spacing = line
    if align is not None: p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.name = font; f.bold = bold; f.italic = italic
    f.color.rgb = RGB(color)
    return p

def runp(tf, parts, size, font=SANS, first=False, before=0, after=0,
         line=None, align=None):
    """paragraph with multiple runs: parts = [(text,color,bold,italic), ...]"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if before: p.space_before = Pt(before)
    if after:  p.space_after = Pt(after)
    if line:   p.line_spacing = line
    if align is not None: p.alignment = align
    for txt, color, bold, *rest in parts:
        it = rest[0] if rest else False
        r = p.add_run(); r.text = txt
        f = r.font; f.size = Pt(size); f.name = font; f.bold = bold
        f.italic = it; f.color.rgb = RGB(color)
    return p

def card(slide, l, t, w, h, fill, radius=0.09):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(l), In(t), In(w), In(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = RGB(fill)
    sp.line.fill.background()
    sp.shadow.inherit = False
    try: sp.adjustments[0] = radius
    except Exception: pass
    return sp

def rule(slide, l, t, w, color=INK, weight=1.25):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, In(l), In(t), In(l+w), In(t))
    ln.line.color.rgb = RGB(color); ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln

def set_bg(slide, color):
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    cSld = slide.element.find(PINK + "cSld")
    for b in cSld.findall(PINK + "bg"):
        cSld.remove(b)
    from lxml import etree as ET
    bg = ET.Element(PINK + "bg"); bgPr = ET.SubElement(bg, PINK + "bgPr")
    sf = ET.SubElement(bgPr, "{%s}solidFill" % A)
    ET.SubElement(sf, "{%s}srgbClr" % A).set("val", color)
    ET.SubElement(bgPr, "{%s}effectLst" % A)
    cSld.insert(0, bg)

def place_image(pic, l, t, w):
    """reposition picture preserving aspect, target width w (in), top t, left l"""
    ar = pic.height / pic.width
    pic.left = In(l); pic.top = In(t); pic.width = In(w); pic.height = In(w*ar)
    return pic

SL = prs.slides

# ================= SLIDE 1 — TITLE =================
s = SL[0]; clear_keep_pics(s); set_bg(s, DARKBG)
_, tf = box(s, 0.9, 1.95, 11, 0.4)
para(tf, "SEMINAR", 11, MUTEDDARK, SANS, bold=True, first=True)
_, tf = box(s, 0.88, 2.45, 11.5, 1.4)
para(tf, "Harness Engineering", 60, ONDARK, SERIF, bold=True, first=True)
rule(s, 0.92, 4.0, 4.2, RULELT, 1.5)
_, tf = box(s, 0.9, 4.25, 10.5, 0.9)
para(tf, "Building and Operating Reliable AI Agents with Claude Code",
     18, MUTEDDARK, SANS, first=True, line=1.3)

# ================= SLIDE 2 — AGENDA (numbered list) =================
s = SL[1]; clear_keep_pics(s)
_, tf = box(s, 0.7, 0.5, 12, 0.8); para(tf, "Agenda", 40, INK, SERIF, bold=True, first=True)
_, tf = box(s, 0.72, 1.35, 12, 0.4)
para(tf, "Six pillars of the harness, then watch them run in two demos.", 15, MUTED, SANS, first=True)
rule(s, 0.7, 1.92, 11.93)
items = [("01","What Is Harness Engineering?"),
         ("02","Context Engineering: The Working Memory Budget"),
         ("03","Orchestration"),
         ("04","Constraints, Guardrails & Safe Autonomy"),
         ("05","Specs, Agent Files & Workflow Design"),
         ("06","Evals & Observability"),
         ("D","Demo")]
y = 2.32; step = 0.625
for num, label in items:
    _, tf = box(s, 0.7, y, 0.9, 0.55, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, num, 22, MUTED, SERIF, bold=True, first=True)
    _, tf = box(s, 1.65, y, 10.8, 0.55, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, label, 17, INK, SANS, bold=False, first=True)
    rule(s, 1.65, y+0.55, 10.95, "D8D3C7", 0.75)
    y += step

# ================= SLIDE 3 — WHY AGENTS FAIL (2x2 cards + footer) =================
s = SL[2]; clear_keep_pics(s)
_, tf = box(s, 0.7, 0.45, 12, 0.8); para(tf, "Why Agents Fail", 40, INK, SERIF, bold=True, first=True)
_, tf = box(s, 0.72, 1.3, 12, 0.4)
para(tf, "The problem is the Harness, not the Model.", 15, MUTED, SANS, first=True)
fails = [("01","Missing Context","Lacks files, docs, or history for informed decisions."),
         ("02","Wrong Tools","Misconfigured or missing tool access — workarounds, errors."),
         ("03","No Guardrails","No scope locks or finish gates — overreach, premature “done”."),
         ("04","No Verification","No automated tests — errors undetected until production.")]
# checker: TL light, TR dark, BL dark, BR light
geom = [(0.7,1.85,5.9,1.92,False),(6.73,1.85,5.9,1.92,True),
        (0.7,3.87,5.9,1.92,True),(6.73,3.87,5.9,1.92,False)]
for (num,head,bodytxt),(l,t,w,h,dark) in zip(fails, geom):
    card(s,l,t,w,h, INK if dark else CARD_L)
    lab = MUTEDDARK if dark else MUTED
    hc  = ONDARK if dark else INK
    bc  = MUTEDDARK if dark else BODY
    _, tf = box(s, l+0.4, t+0.28, w-0.8, h-0.5)
    para(tf, "PROFILE "+num if False else num, 11, lab, SANS, bold=True, first=True)
    para(tf, head, 24, hc, SERIF, bold=True, before=4)
    para(tf, bodytxt, 13, bc, SANS, before=8, line=1.3)
# footer callout
_, tf = box(s, 0.7, 6.18, 11.93, 0.6, anchor=MSO_ANCHOR.MIDDLE)
runp(tf, [("Key Insight:  ", INK, True), ("Most failures live in the Harness — not the Model.", MUTED, False)],
     15, SANS, first=True, align=PP_ALIGN.CENTER)

# ================= SECTION DIVIDERS 4,6,8,10,12,14 =================
dividers = {
 4:("01","What a Harness Actually is?","The model supplies reasoning. The harness supplies discipline."),
 6:("02","Context Engineering","The Working Memory Budget — the context window is a finite resource."),
 8:("03","Orchestration","One agent doing it all = full context, hallucination. Split orchestrator & sub-agents."),
 10:("04","Constraints, Guardrails & Safe Autonomy","Agents fail from missing scope discipline, not missing reasoning."),
 12:("05","Specs, Agent Files & Workflow Design","Instruction & context layer — load progressively, not the whole encyclopedia."),
 14:("06","Evals & Observability","Agent confidence ≠ evidence. Data is evidence: traces, citations, evals."),
}
for n,(num,title,sub) in dividers.items():
    s = SL[n-1]; clear_keep_pics(s); set_bg(s, DARKBG)
    _, tf = box(s, 0.85, 2.05, 3.0, 2.0, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, num, 130, ONDARK, SERIF, bold=True, first=True)
    _, tf = box(s, 3.45, 2.35, 9.0, 1.9, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, 38, ONDARK, SERIF, bold=True, first=True, line=1.05)
    para(tf, sub, 18, MUTEDDARK, SANS, before=14, line=1.3)
    # thin vertical accent between number and title
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, In(3.3), In(2.5), In(3.3), In(4.0))
    ln.line.color.rgb = RGB(RULELT); ln.line.width = Pt(1.5); ln.shadow.inherit=False

# ================= CONTENT SLIDES 5,7,9,11,13,15 (title + text card + image) =================
content = {
 5:("What a Harness Actually is?","Strong model fails here, succeeds elsewhere — harness, not model."),
 7:("Context Engineering: Working Memory Budget","Manage the budget so the agent keeps decisions through compaction."),
 9:("Orchestration & Sub-Agents","Control plane runs the query loop; execution plane runs only permitted tools."),
 11:("Constraints, Guardrails & Safe Autonomy","Scope lock · finish gate · circuit breaker · error as the main path."),
 13:("Specs, Agent Files & Workflow Design","AGENTS.md = index · feature list = primitive · init precedes action."),
 15:("Evals & Observability","Validate behavior with evidence and traces, not predictions."),
}
for n,(title,sub) in content.items():
    s = SL[n-1]
    pic = get_pic(s)
    clear_keep_pics(s)
    # bullets = the big text box (3rd text box on these slides) -> find longest list
    paras = slide_paras(n)
    bullets = max(paras, key=len)  # the multi-bullet box
    bullets = [b.lstrip("—").strip().lstrip("-").strip() for b in bullets]
    _, tf = box(s, 0.7, 0.5, 12.3, 0.7); para(tf, title, 29, INK, SERIF, bold=True, first=True)
    _, tf = box(s, 0.72, 1.18, 12.3, 0.4); para(tf, sub, 15, MUTED, SANS, first=True)
    rule(s, 0.7, 1.72, 11.93)
    # left text card
    card(s, 0.7, 1.98, 6.45, 4.75, CARD_L)
    _, tf = box(s, 1.02, 2.3, 5.85, 4.2)
    for i,b in enumerate(bullets):
        runp(tf, [("—  ", MUTED, True), (b, BODY, False)], 12.5, SANS,
             first=(i==0), after=9, line=1.18)
    # right image
    if pic is not None:
        # fit within right column: x 7.4..12.85 (w<=5.45), y from 2.1, max h 4.5
        ar = pic.height/pic.width
        w = 5.45; h = w*ar
        if h > 4.55:
            h = 4.55; w = h/ar
        x = 7.45 + (5.45 - w)/2
        pic.left = In(x); pic.top = In(2.15); pic.width = In(w); pic.height = In(h)

# ================= SLIDE 16 — architecture image =================
s = SL[15]
pic = get_pic(s); clear_keep_pics(s)
_, tf = box(s, 0.7, 0.5, 12.3, 0.7)
para(tf, "Claude Code Harness Architecture", 29, INK, SERIF, bold=True, first=True)
_, tf = box(s, 0.72, 1.18, 12.3, 0.5)
para(tf, "Six subsystems — Cache, Security, Tool Orchestration, Memory, State, Multi-Agent — wrapped around the LLM.",
     15, MUTED, SANS, first=True, line=1.25)
rule(s, 0.7, 1.78, 11.93)
if pic is not None:
    ar = pic.height/pic.width
    h = 4.55; w = h/ar
    pic.left = In((13.33-w)/2); pic.top = In(2.05); pic.width = In(w); pic.height = In(h)
_, tf = box(s, 0.7, 6.92, 12, 0.3)
para(tf, "Source: Diagram synthesized from leaked documentation on the Claude Code Agent Harness architecture.",
     11, MUTED, SANS, first=True, align=PP_ALIGN.CENTER)

# ================= SLIDE 17 — LIVE DEMO (statement) =================
s = SL[16]; clear_keep_pics(s); set_bg(s, DARKBG)
rule(s, 0.92, 2.62, 2.2, RULELT, 2.0)
_, tf = box(s, 0.9, 2.78, 11.5, 1.5)
para(tf, "LIVE DEMO", 72, ONDARK, SERIF, bold=True, first=True)
_, tf = box(s, 0.92, 4.35, 11, 0.6)
para(tf, "Theory, now for real — two harnesses built with Claude Code.", 18, MUTEDDARK, SANS, first=True)

# ================= SLIDE 21 — KEY TAKEAWAYS (list + dark footer bar) =================
s = SL[20]; clear_keep_pics(s)
_, tf = box(s, 0.7, 0.5, 12, 0.8); para(tf, "Key Takeaways", 40, INK, SERIF, bold=True, first=True)
rule(s, 0.7, 1.5, 11.93)
takeaways = [
 "Model = reasoning · Harness = discipline + observable execution. Failures live in the harness.",
 "Context = budget: progressive disclosure · repo as system of record · compaction-aware.",
 "Orchestrator holds summaries; sub-agents return summaries — context isolation = scale.",
 "Done = evidence (test/trace), not a promise · Guardrails in the control plane.",
 "Provenance citations + golden-question evals = trust you can MEASURE.",
]
y = 1.82; step = 0.86
for i,t in enumerate(takeaways,1):
    _, tf = box(s, 0.7, y, 0.85, 0.7, anchor=MSO_ANCHOR.TOP)
    para(tf, f"{i:02d}", 20, MUTED, SERIF, bold=True, first=True)
    _, tf = box(s, 1.6, y, 11.0, 0.78, anchor=MSO_ANCHOR.TOP)
    para(tf, t, 16, INK, SANS, first=True, line=1.2)
    if i < len(takeaways):
        rule(s, 1.6, y+step-0.16, 11.0, "D8D3C7", 0.75)
    y += step
# dark footer callout bar
card(s, 0.7, 6.42, 11.93, 0.72, INK, radius=0.18)
_, tf = box(s, 0.9, 6.42, 11.5, 0.72, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Claude Code is already a harness — operate it with discipline.",
     15, ONDARK, SANS, bold=True, first=True, align=PP_ALIGN.CENTER)

# ================= SLIDE 22 — THANK YOU (closing statement) =================
s = SL[21]; clear_keep_pics(s); set_bg(s, DARKBG)
_, tf = box(s, 0.9, 2.2, 11.5, 1.1)
para(tf, "Thank You — Q&A", 54, ONDARK, SERIF, bold=True, first=True)
_, tf = box(s, 0.92, 3.35, 11.5, 0.5)
para(tf, "Harness Engineering: Building and Operating Reliable AI Agents with Claude Code",
     16, MUTEDDARK, SANS, first=True)
rule(s, 0.92, 4.05, 4.2, RULELT, 1.5)
_, tf = box(s, 0.9, 4.3, 11.5, 1.0)
para(tf, "“A smart model without a harness is just an intern with root access.”",
     20, ONDARK, SERIF, italic=True, first=True, line=1.25)

# diagram slides 18,19,20 are kept as-is; just drop the thin dark left-edge accent bar
for idx in (17, 18, 19):
    s = SL[idx]
    for sh in list(s.shapes):
        if sh.shape_type == 13:
            continue
        if (sh.left is not None and sh.left < In(0.06)
                and sh.width is not None and sh.width < In(0.3)
                and sh.height is not None and sh.height > In(6)):
            sh._element.getparent().remove(sh._element)

prs.save(OUT)
print("rebuilt -> ", OUT)
